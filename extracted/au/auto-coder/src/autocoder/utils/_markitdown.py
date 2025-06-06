# type: ignore
import base64
import binascii
import copy
import html
import json
import mimetypes
import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import uuid
from typing import Any, Dict, List, Optional, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LAParams, LTImage, LTFigure
from pdfminer.pdfdevice import PDFDevice
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
import pptx
from pdfminer.image import ImageWriter
import time

import numpy as np
from PIL import Image

# 新增导入
from autocoder.rag.loaders.filter_utils import FilterRuleManager
from autocoder.rag.loaders.image_loader import ImageLoader

# File-format detection
import puremagic
import requests
from bs4 import BeautifulSoup
from loguru import logger

# Optional Transcription support
try:
    import pydub
    import speech_recognition as sr

    IS_AUDIO_TRANSCRIPTION_CAPABLE = True
except ModuleNotFoundError:
    pass

# Optional YouTube transcription support
try:
    from youtube_transcript_api import YouTubeTranscriptApi

    IS_YOUTUBE_TRANSCRIPT_CAPABLE = True
except ModuleNotFoundError:
    pass


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get(
            "heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in [
                    "http",
                    "https",
                    "file",
                ]:
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(
                    parsed_url._replace(path=quote(unquote(parsed_url.path)))
                )  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        try:
            # 设置递归深度限制，避免复杂文档导致的递归错误
            import sys
            original_limit = sys.getrecursionlimit()
            try:
                # 增加递归深度限制
                sys.setrecursionlimit(10000)  # 设置更高的递归限制
                return super().convert_soup(soup)  # type: ignore
            finally:
                # 恢复原始递归深度限制
                sys.setrecursionlimit(original_limit)
        except RecursionError:
            # 处理递归错误，尝试简化处理
            logger.warning("RecursionError in convert_soup, falling back to simplified conversion")
            # 返回简化的文本内容
            return self._simplified_convert(soup)

    def _simplified_convert(self, soup: Any) -> str:
        """简化的转换方法，用于处理复杂文档时的回退方案"""
        # 提取纯文本内容
        text = soup.get_text(separator="\n", strip=True)
        # 基本清理
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Guess the content type from any file extension that might be around
        content_type, _ = mimetypes.guess_type(
            "__placeholder" + kwargs.get("file_extension", "")
        )

        # Only accept text files
        if content_type is None:
            return None
        elif "text/" not in content_type.lower():
            return None

        text_content = ""
        with open(local_path, "rt", encoding="utf-8") as fh:
            text_content = fh.read()
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        try:
            # Print only the main content
            body_elm = soup.find("body")
            webpage_text = ""
            if body_elm:
                webpage_text = _CustomMarkdownify().convert_soup(body_elm)
            else:
                webpage_text = _CustomMarkdownify().convert_soup(soup)

            assert isinstance(webpage_text, str)

            return DocumentConverterResult(
                title=None if soup.title is None else soup.title.string,
                text_content=webpage_text,
            )
        except Exception as e:
            # 如果转换过程中出现任何错误，尝试使用简化的方法提取文本
            logger.warning(f"Error in HTML conversion: {str(e)}. Falling back to simplified text extraction.")
            try:
                # 简化的文本提取
                text = soup.get_text(separator="\n", strip=True)
                # 基本清理
                text = re.sub(r'\s+', ' ', text)
                text = re.sub(r'\n{3,}', '\n\n', text)
                
                return DocumentConverterResult(
                    title=None if soup.title is None else soup.title.string,
                    text_content=text,
                )
            except Exception as inner_e:
                # 如果简化提取也失败，记录错误并返回空结果
                logger.error(f"Failed to extract text with simplified method: {str(inner_e)}")
                return DocumentConverterResult(
                    title=None,
                    text_content=f"[文档转换失败] 无法提取内容: {str(e)}",
                )


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not Wikipedia
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        main_title = None if soup.title is None else soup.title.string

        if body_elm:
            # What's the title
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string  # type: ignore
                assert isinstance(main_title, str)

            # Convert the page
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(
                body_elm
            )
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=main_title,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        assert soup.title is not None and soup.title.string is not None
        metadata: Dict[str, str] = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start: obj_end + 1])
                        attrdesc = self._findKey(
                            data, "attributedDescriptionBodyText"
                        )  # type: ignore
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(
            metadata, ["title", "og:title", "name"])  # type: ignore
        assert isinstance(title, str)

        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])  # type: ignore
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])  # type: ignore
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])  # type: ignore
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(
            metadata, ["description", "og:description"]
        )  # type: ignore
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        if IS_YOUTUBE_TRANSCRIPT_CAPABLE:
            transcript_text = ""
            parsed_url = urlparse(url)  # type: ignore
            params = parse_qs(parsed_url.query)  # type: ignore
            if "v" in params:
                assert isinstance(params["v"][0], str)
                video_id = str(params["v"][0])
                try:
                    # Must be a single transcript.
                    transcript = YouTubeTranscriptApi.get_transcript(
                        video_id
                    )  # type: ignore
                    transcript_text = " ".join(
                        [part["text"] for part in transcript]
                    )  # type: ignore
                    # Alternative formatting:
                    # formatter = TextFormatter()
                    # formatter.format_transcript(transcript)
                except Exception:
                    pass
            if transcript_text:
                webpage_text += f"\n### Transcript\n{transcript_text}\n"

        title = title if title else soup.title.string
        assert isinstance(title, str)

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
        )

    def _get(
        self,
        metadata: Dict[str, str],
        keys: List[str],
        default: Union[str, None] = None,
    ) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    # TODO: Fix json type
    def _findKey(self, json: Any, key: str) -> Union[str, None]:
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None


class BingSerpConverter(DocumentConverter):
    """
    Handle Bing results pages (only the organic search results).
    NOTE: It is better to use the Bing API
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a Bing SERP
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https://www\.bing\.com/search\?q=", url):
            return None

        # Parse the query parameters
        parsed_params = parse_qs(urlparse(url).query)
        query = parsed_params.get("q", [""])[0]

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Clean up some formatting
        for tptt in soup.find_all(class_="tptt"):
            if hasattr(tptt, "string") and tptt.string:
                tptt.string += " "
        for slug in soup.find_all(class_="algoSlug_icon"):
            slug.extract()

        # Parse the algorithmic results
        _markdownify = _CustomMarkdownify()
        results = list()
        for result in soup.find_all(class_="b_algo"):
            # Rewrite redirect urls
            for a in result.find_all("a", href=True):
                parsed_href = urlparse(a["href"])
                qs = parse_qs(parsed_href.query)

                # The destination is contained in the u parameter,
                # but appears to be base64 encoded, with some prefix
                if "u" in qs:
                    u = (
                        qs["u"][0][2:].strip() + "=="
                    )  # Python 3 doesn't care about extra padding

                    try:
                        # RFC 4648 / Base64URL" variant, which uses "-" and "_"
                        a["href"] = base64.b64decode(
                            u, altchars="-_").decode("utf-8")
                    except UnicodeDecodeError:
                        pass
                    except binascii.Error:
                        pass

            # Convert to markdown
            md_result = _markdownify.convert_soup(result).strip()
            lines = [line.strip() for line in re.split(r"\n+", md_result)]
            results.append(
                "\n".join([line for line in lines if len(line) > 0]))

        webpage_text = (
            f"## A Bing search for '{query}' found the following results:\n\n"
            + "\n\n".join(results)
        )

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown with support for extracting and including images.
    """

    def __init__(self, llm=None, product_mode="lite"):
        super().__init__()
        self.llm = llm
        self.product_mode = product_mode

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None        
        image_output_dir = None
        if kwargs.get("image_output_dir", None):
            image_output_dir = kwargs.get("image_output_dir")
        else:
            # Create output directory for images if it doesn't exist
            image_output_dir = os.path.join(
                os.path.dirname(local_path), "_images", os.path.basename(
                    local_path).replace(" ", "_")
            )
        os.makedirs(image_output_dir, exist_ok=True)

        text_content = []
        image_count = 0

        # Open and process PDF
        with open(local_path, "rb") as file:
            # Create PDF parser and document
            parser = PDFParser(file)
            document = PDFDocument(parser)
            rsrcmgr = PDFResourceManager()
            laparams = LAParams()
            device = PDFPageAggregator(rsrcmgr, laparams=laparams)
            interpreter = PDFPageInterpreter(rsrcmgr, device)            

            # Process each page
            for page in PDFPage.create_pages(document):
                interpreter.process_page(page)
                layout = device.get_result()                

                # Extract text and images from the page
                page_content = self._process_layout(
                    layout, image_output_dir, image_count
                )
                
                text_content.extend(page_content)
                image_count += len([c for c in page_content if c.startswith("![Image")])

        return DocumentConverterResult(
            title=None,
            text_content="\n".join(text_content),
        )

    def _process_layout(
        self, layout, image_output_dir: str, image_count: int
    ) -> List[str]:
        """Process the layout of a PDF page, extracting both text and images."""
        content = [] 
        local_image_count = image_count       
        for lt_obj in layout:
            # Handle images
            if isinstance(lt_obj, LTImage) or (
                isinstance(lt_obj, LTFigure) and lt_obj.name.startswith("Im")
            ):                
                image_data = None
                image_meta = {}
                # 添加uuid后缀避免文件名冲突
                unique_id = str(uuid.uuid4())[:8]
                image_path = os.path.join(
                    image_output_dir, f"image_{local_image_count}_{unique_id}.png")

                if hasattr(lt_obj, "stream"):
                    image_data = lt_obj.stream.get_data()
                    image_meta = lt_obj.stream.attrs
                elif hasattr(lt_obj, "filter"):
                    image_data = lt_obj.filter

                if image_data:
                    if isinstance(lt_obj, LTImage):
                        iw = ImageWriter(image_output_dir)
                        name = iw.export_image(lt_obj)
                        suffix = os.path.splitext(name)[1]
                        temp_path = os.path.join(image_output_dir, name)
                        # 添加uuid后缀避免文件名冲突
                        unique_id = str(uuid.uuid4())[:8]
                        image_path = os.path.join(
                            image_output_dir, f"image_{local_image_count}_{unique_id}{suffix}")
                        os.rename(temp_path, image_path)
                        content.append(f"![Image {local_image_count}]({image_path})")
                        # ===== 修改：通过FilterRuleManager单例实例判断是否需要解析图片
                        v = try_parse_image(image_path,self.llm)
                        if v:
                            content.append("<image_content>")
                            content.append(v)
                            content.append("</image_content>")
                        # =====
                        local_image_count += 1
                        continue
                    try:
                        # Try to handle raw pixel data
                        if "BitsPerComponent" in image_meta:
                            width = image_meta["Width"]
                            height = image_meta["Height"]
                            bits = image_meta["BitsPerComponent"]
                            colorspace = image_meta["ColorSpace"].name
                            new_image_data = np.frombuffer(
                                image_data, dtype=np.uint8)
                            # Normalize to 8-bit if necessary
                            if bits != 8:
                                max_val = (1 << bits) - 1
                                new_image_data = (
                                    new_image_data.astype(
                                        "float32") * 255 / max_val
                                ).astype("uint8")

                            if colorspace == "DeviceRGB":
                                new_image_data = new_image_data.reshape(
                                    (height, width, 3)
                                )
                                img = Image.fromarray(new_image_data, "RGB")
                                # 添加uuid后缀避免文件名冲突
                                unique_id = str(uuid.uuid4())[:8]
                                image_path = os.path.join(
                                    image_output_dir, f"image_{local_image_count}_{unique_id}.png")
                                img.save(image_path)
                                content.append(
                                    f"![Image {local_image_count}]({image_path})\n"
                                )
                                v = try_parse_image(image_path,self.llm)
                                if v:
                                    content.append("<image_content>")
                                    content.append(v)
                                    content.append("</image_content>")
                                local_image_count += 1
                                continue
                            elif colorspace == "DeviceGray":
                                new_image_data = new_image_data.reshape(
                                    (height, width))
                                img = Image.fromarray(new_image_data, "L")
                                # 添加uuid后缀避免文件名冲突
                                unique_id = str(uuid.uuid4())[:8]
                                image_path = os.path.join(
                                    image_output_dir, f"image_{local_image_count}_{unique_id}.png")
                                img.save(image_path)
                                content.append(
                                    f"![Image {local_image_count}]({image_path})\n"
                                )
                                v = try_parse_image(image_path,self.llm)
                                if v:
                                    content.append("<image_content>")
                                    content.append(v)
                                    content.append("</image_content>")
                                local_image_count += 1
                                continue
                    except Exception as e:
                        print(
                            f"Error extracting image: {e} fallback to writing original data"
                        )

                    # 添加uuid后缀避免文件名冲突
                    unique_id = str(uuid.uuid4())[:8]
                    image_path = os.path.join(
                        image_output_dir, f"image_{local_image_count}_{unique_id}.png")
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_data)

                    content.append(f"![Image {local_image_count}]({image_path})\n")
                    # ===== 新增：图片解析
                    v = try_parse_image(image_path,self.llm)
                    if v:
                        content.append("<image_content>")
                        content.append(v)
                        content.append("</image_content>")
                    local_image_count += 1

            # Handle text
            if hasattr(lt_obj, "get_text"):
                text = lt_obj.get_text().strip()
                if text:
                    content.append(text)

            # Recursively process nested layouts
            elif hasattr(lt_obj, "_objs"):
                content.extend(
                    self._process_layout(
                        lt_obj._objs, image_output_dir, image_count)
                )

        return content


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def __init__(self):
        self._image_counter = 0
        super().__init__()

    def _save_image(self, image, output_dir: str) -> str:
        """
        保存图片并返回相对路径，使用递增的计数器来命名文件
        """
        # 获取图片内容和格式
        image_format = image.content_type.split(
            '/')[-1] if image.content_type else 'png'

        # 增加计数器并生成文件名
        self._image_counter += 1
        # 添加uuid后缀避免文件名冲突
        unique_id = str(uuid.uuid4())[:8]
        image_filename = f"image_{self._image_counter}_{unique_id}.{image_format}"

        # 保存图片
        image_path = os.path.join(output_dir, image_filename)
        with image.open() as image_content, open(image_path, 'wb') as f:
            f.write(image_content.read())

        return image_path

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        # 设置图片输出目录
        image_output_dir = None
        if kwargs.get("image_output_dir", None):
            image_output_dir = kwargs.get("image_output_dir")
        else:
            # Create output directory for images if it doesn't exist
            image_output_dir = os.path.join(os.path.dirname(
                local_path), "_images", os.path.basename(local_path).replace(" ", "_"))
        os.makedirs(image_output_dir, exist_ok=True)

        result = None
        with open(local_path, "rb") as docx_file:
            # 配置图片转换器
            def transform_image(image):
                return {
                    "src": self._save_image(image, image_output_dir),
                    "alt": image.alt_text if image.alt_text else f"Image {self._image_counter}"
                }

            # 进行转换
            result = mammoth.convert_to_html(
                docx_file,
                convert_image=mammoth.images.inline(transform_image)
            )
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(
                html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get(
                            "descr", "")
                    except Exception:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += (
                        "\n!["
                        + (alt_text if alt_text else shape.name)
                        + "]("
                        + filename
                        + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + \
                                    html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + \
                                    html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" +
                        self._convert(html_table).text_content.strip() + "\n"
                    )

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run(
                    [exiftool, "-json", local_path], capture_output=True, text=True
                ).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            try:
                transcript = self._transcribe_audio(local_path)
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += (
                    "\n\n### Audio Transcript:\nError. Could not transcribe this audio."
                )

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".mp3":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            handle, temp_path = tempfile.mkstemp(suffix=".wav")
            os.close(handle)
            try:
                sound = pydub.AudioSegment.from_mp3(local_path)
                sound.export(temp_path, format="wav")

                _args = dict()
                _args.update(kwargs)
                _args["file_extension"] = ".wav"

                try:
                    transcript = super()._transcribe_audio(temp_path).strip()
                    md_content += "\n\n### Audio Transcript:\n" + (
                        "[No speech detected]" if transcript == "" else transcript
                    )
                except Exception:
                    md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

            finally:
                os.unlink(temp_path)

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an mlm_client is configured).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        mlm_client = kwargs.get("mlm_client")
        mlm_model = kwargs.get("mlm_model")
        if mlm_client is not None and mlm_model is not None:
            md_content += (
                "\n# Description:\n"
                + self._get_mlm_description(
                    local_path,
                    extension,
                    mlm_client,
                    mlm_model,
                    prompt=kwargs.get("mlm_prompt"),
                ).strip()
                + "\n"
            )

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _get_mlm_description(self, local_path, extension, client, model, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        sys.stderr.write(f"MLM Prompt:\n{prompt}\n")

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                ],
            }
        ]

        response = client.chat.completions.create(
            model=model, messages=messages)
        return response.choices[0].message.content


class FileConversionException(BaseException):
    pass


class UnsupportedFormatException(BaseException):
    pass


class MarkItDown:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[Any] = None,
        llm: Optional[Any] = None,
        product_mode: Optional[str] = None,
    ):
        # 初始化FilterRuleManager单例实例
        self._filter_rule_manager = FilterRuleManager.get_instance()
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client
        self._mlm_model = mlm_model

        # 新增：保存llm和product_mode
        self._llm = llm
        self._product_mode = product_mode

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(BingSerpConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(PdfConverter(llm,product_mode))

    def convert(
        self, source: Union[str, requests.Response], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """        
        # Local path or url
        if isinstance(source, str):
            if (
                source.startswith("http://")
                or source.startswith("https://")
                or source.startswith("file://")
            ):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(
        self, path: str, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)

        if not extensions:
            for g in self._guess_ext_magic(path):
                self._append_ext(extensions, g)

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(
        self, stream: Any, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(
        self, url: str, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        response = self._requests_session.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, url=response.url)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(
        self, local_path: str, extensions: List[Union[str, None]], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""
        res = None
        for ext in extensions + [None]:  # Try last with no extension
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                # Copy any additional global options
                if "mlm_client" not in _kwargs and self._mlm_client is not None:
                    _kwargs["mlm_client"] = self._mlm_client

                if "mlm_model" not in _kwargs and self._mlm_model is not None:
                    _kwargs["mlm_model"] = self._mlm_model

                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join(
                        [line.rstrip()
                         for line in re.split(r"\r?\n", res.text_content)]
                    )
                    res.text_content = re.sub(
                        r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            extensions = list()
            for g in guesses:
                ext = g.extension.strip()
                if len(ext) > 0:
                    if not ext.startswith("."):
                        ext = "." + ext
                    if ext not in extensions:
                        extensions.append(ext)
            return extensions
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return []

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)


def try_parse_image(image_path: str, llm=None):
    """
    根据FilterRuleManager单例实例判断是否需要解析图片，如果需要则调用ImageLoader.image_to_markdown。
    解析失败会自动捕获异常。
    """
    import uuid
    start_time = time.time()
    req_id = str(uuid.uuid4())[:8]
    logger.info(f"\n==== [try_parse_image] START | req_id={req_id} ====")
    logger.info(f"[try_parse_image][{req_id}] image_path: {image_path}, llm: {llm}")
    if FilterRuleManager.get_instance().should_parse_image(image_path):
        logger.info(f"[try_parse_image][{req_id}] should_parse_image=True, start parsing...")
        try:
            v = ImageLoader.image_to_markdown(image_path, llm=llm, engine="paddle")
            logger.info(f"[try_parse_image][{req_id}] image_to_markdown result: {str(v)[:200]}")
            
            if not v:
                return ""
            
            if llm:
                v = ImageLoader.format_table_in_content(v, llm)
                logger.info(f"[try_parse_image][{req_id}] format_table_in_content result: {str(v)[:200]}")
            elapsed = time.time() - start_time
            logger.info(f"[try_parse_image][{req_id}] SUCCESS | execution time: {elapsed:.3f} seconds")
            logger.info(f"==== [try_parse_image] END | req_id={req_id} ====")
            return v
        except Exception as e:
            elapsed = time.time() - start_time
            logger.error(f"[try_parse_image][{req_id}] EXCEPTION | execution time: {elapsed:.3f} seconds | image_path: {image_path} | llm: {llm}")
            logger.exception(e)
            logger.info(f"==== [try_parse_image] END (EXCEPTION) | req_id={req_id} ====")
            return ""
    else:
        logger.info(f"[try_parse_image][{req_id}] should_parse_image=False, skip parsing.")
        logger.info(f"==== [try_parse_image] END (SKIP) | req_id={req_id} ====")
        return ""

